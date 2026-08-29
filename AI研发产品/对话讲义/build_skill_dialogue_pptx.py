# -*- coding: utf-8 -*-
"""
生成讲义式 PPT：从「找自动化 skill」到「API / Auto / Codex / 模型」探讨。
风格：大白话开场总结 + 完整对话回放，说话人清晰可辨。
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent
OUT = ROOT / "从Skill探讨到AI对比_讲义.pptx"
ART = Path("/opt/cursor/artifacts/skill-dialogue-ppt")

# Palette: ink-on-paper lecture, not purple-AI default
BG = RGBColor(0xF7, 0xF3, 0xEA)          # warm paper
INK = RGBColor(0x1C, 0x24, 0x2E)          # near-black ink
MUTED = RGBColor(0x5B, 0x64, 0x72)        # secondary text
LINE = RGBColor(0xD7, 0xCF, 0xC0)         # hairline
USER_BG = RGBColor(0xE8, 0xF2, 0xEF)      # soft teal wash
USER_ACCENT = RGBColor(0x1F, 0x6F, 0x64)  # teal
AI_BG = RGBColor(0xF3, 0xEB, 0xE0)         # warm sand
AI_ACCENT = RGBColor(0xA6, 0x4B, 0x2A)     # terracotta / brick
TITLE_BG = RGBColor(0x1C, 0x24, 0x2E)
TITLE_FG = RGBColor(0xF7, 0xF3, 0xEA)
ACCENT = RGBColor(0xC4, 0x5C, 0x26)

FONT = "Microsoft YaHei"
FONT_FALLBACK = "WenQuanYi Micro Hei"


def _set_run_font(run, size_pt: float, bold: bool = False, color: RGBColor | None = None, name: str = FONT):
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.name = name
    if color is not None:
        run.font.color.rgb = color
    # East Asian font hint for Chinese
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", name)


def _fill(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _add_textbox(slide, left, top, width, height, text, *, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT, font=FONT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_run_font(run, size, bold=bold, color=color, name=font)
    return box


def _add_paras(slide, left, top, width, height, lines: list[tuple[str, dict]], *, font=FONT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for text, opt in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = opt.get("align", PP_ALIGN.LEFT)
        p.space_after = Pt(opt.get("after", 8))
        run = p.add_run()
        run.text = text
        _set_run_font(
            run,
            opt.get("size", 16),
            bold=opt.get("bold", False),
            color=opt.get("color", INK),
            name=font,
        )
    return box


def blank_slide(prs: Presentation, dark: bool = False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    _fill(bg, TITLE_BG if dark else BG)
    return slide


def add_footer(slide, prs, page: str, dark: bool = False):
    color = RGBColor(0x9A, 0xA3, 0xB0) if dark else MUTED
    _add_textbox(
        slide,
        Inches(0.6),
        Inches(7.05),
        Inches(10.5),
        Inches(0.3),
        f"Skill 探讨讲义  ·  {page}",
        size=11,
        color=color,
    )


def cover(prs):
    slide = blank_slide(prs, dark=True)
    # accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.7), Inches(1.2), Inches(0.12))
    _fill(bar, ACCENT)
    _add_textbox(slide, Inches(0.7), Inches(2.0), Inches(11), Inches(1.2),
                 "从「找一个自动化 Skill」开始", size=34, bold=True, color=TITLE_FG)
    _add_textbox(slide, Inches(0.7), Inches(3.2), Inches(11), Inches(1.0),
                 "到弄明白：说明书、API、模型、平台，各自干什么", size=22, color=RGBColor(0xD2, 0xC6, 0xB4))
    _add_textbox(slide, Inches(0.7), Inches(5.2), Inches(11), Inches(0.8),
                 "一场大白话对话讲义  ·  你问我答  ·  先总结方向，再回放全文", size=15, color=RGBColor(0x9A, 0xA3, 0xB0))
    add_footer(slide, prs, "封面", dark=True)


def section(prs, eyebrow: str, title: str, subtitle: str = ""):
    slide = blank_slide(prs, dark=True)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(2.3), Inches(0.9), Inches(0.1))
    _fill(bar, ACCENT)
    _add_textbox(slide, Inches(0.7), Inches(1.6), Inches(11), Inches(0.4), eyebrow, size=14, color=ACCENT, bold=True)
    _add_textbox(slide, Inches(0.7), Inches(2.55), Inches(11), Inches(1.2), title, size=30, bold=True, color=TITLE_FG)
    if subtitle:
        _add_textbox(slide, Inches(0.7), Inches(4.0), Inches(11), Inches(1.2), subtitle, size=16, color=RGBColor(0xC8, 0xBD, 0xAD))
    add_footer(slide, prs, "分节", dark=True)


def summary_why(prs):
    slide = blank_slide(prs)
    _add_textbox(slide, Inches(0.7), Inches(0.45), Inches(11), Inches(0.4), "开场总结 · 为什么会聊到后面那些", size=13, bold=True, color=ACCENT)
    _add_textbox(slide, Inches(0.7), Inches(0.85), Inches(11), Inches(0.6), "你不是上来就问模型榜单的", size=26, bold=True, color=INK)
    lines = [
        ("你真正想办的事很具体：", {"size": 16, "bold": True, "after": 10}),
        ("想找能画 CAD、自动处理文档、管邮件、自动回微信的 skill。", {"size": 16, "after": 8}),
        ("结果发现：当前环境里并没有这种现成技能包可“搜索安装”。", {"size": 16, "after": 8}),
        ("所以对话自然拐到了——", {"size": 16, "bold": True, "after": 8}),
        ("skill 到底是不是能力本身？要不要自己开发？Codex 是不是更多？", {"size": 15, "color": MUTED, "after": 6}),
        ("API 又是啥？Auto 和 GPT、Claude 差在哪？阀门定位器要不要上 Opus？", {"size": 15, "color": MUTED, "after": 6}),
        ("后面所有关于 AI 的探讨，都是被这件事逼出来的方向澄清。", {"size": 16, "bold": True, "after": 4}),
    ]
    _add_paras(slide, Inches(0.7), Inches(1.7), Inches(11.2), Inches(5.0), lines)
    add_footer(slide, prs, "01")


def summary_plain(prs):
    slide = blank_slide(prs)
    _add_textbox(slide, Inches(0.7), Inches(0.4), Inches(11), Inches(0.35), "开场总结 · 大白话版", size=13, bold=True, color=ACCENT)
    _add_textbox(slide, Inches(0.7), Inches(0.8), Inches(11), Inches(0.5), "先把几个容易混的词说清楚", size=24, bold=True)

    cards = [
        ("Skill", "说明书 / SOP", "教 AI 遇到这类事先按什么步骤做。\n本身不等于已经连上微信或 CAD。"),
        ("API", "插座 / 接口", "让程序真的去读邮箱、发通知、驱动软件。\n没它，很多自动化只能停留在“会说”。"),
        ("模型", "脑子", "Claude / GPT / Kimi / Grok……\n负责想、写、判断。"),
        ("平台", "身体 + 工具箱", "Cursor Auto / Codex……\n决定看得见项目、改不改得了文件、接不接得上工具。"),
    ]
    for i, (k, sub, body) in enumerate(cards):
        col = i % 2
        row = i // 2
        left = Inches(0.7 + col * 6.1)
        top = Inches(1.55 + row * 2.45)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.8), Inches(2.2))
        _fill(card, RGBColor(0xFF, 0xFC, 0xF7))
        card.line.color.rgb = LINE
        _add_textbox(slide, left + Inches(0.25), top + Inches(0.2), Inches(5.2), Inches(0.4), k, size=18, bold=True, color=USER_ACCENT)
        _add_textbox(slide, left + Inches(0.25), top + Inches(0.55), Inches(5.2), Inches(0.35), sub, size=13, bold=True, color=AI_ACCENT)
        _add_textbox(slide, left + Inches(0.25), top + Inches(0.95), Inches(5.2), Inches(1.0), body, size=13, color=INK)
    add_footer(slide, prs, "02")


def summary_verdict(prs):
    slide = blank_slide(prs)
    _add_textbox(slide, Inches(0.7), Inches(0.45), Inches(11), Inches(0.35), "开场总结 · 后来达成的共识", size=13, bold=True, color=ACCENT)
    _add_textbox(slide, Inches(0.7), Inches(0.9), Inches(11), Inches(0.55), "不是换个软件名字就自动会干活", size=24, bold=True)
    bullets = [
        ("1. 现成 skill 少，不等于装配能力差。", "Auto 能写代码、接 API、搭流程；差的是现成说明书多不多。"),
        ("2. Codex 往往生态更整，但也不等于装完就有微信/CAD。", "还是要 skill + API + 代码一起配。"),
        ("3. 模型像发动机，平台像整车。", "榜单多测脑子；真实项目还看工具、权限、流程。"),
        ("4. 你喂 Blender 出 FY301 视频，本质是：工具补上了 + 流程沉淀成 skill。", "装软件偏工具；一步步纠正偏说明书。"),
        ("5. 孔板要“像 CAD 画的”，得操作 CAD / 用母版，不是只靠 ezdxf 吐文件。", "文丘里图是观感样板，不是让你画文丘里。"),
    ]
    y = 1.65
    for title, body in bullets:
        _add_textbox(slide, Inches(0.7), Inches(y), Inches(11.2), Inches(0.35), title, size=15, bold=True, color=INK)
        _add_textbox(slide, Inches(0.9), Inches(y + 0.32), Inches(11), Inches(0.4), body, size=13, color=MUTED)
        y += 0.95
    add_footer(slide, prs, "03")


def legend(prs):
    slide = blank_slide(prs)
    _add_textbox(slide, Inches(0.7), Inches(0.5), Inches(11), Inches(0.4), "怎么读后面的对话回放", size=24, bold=True)
    # user card
    u = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.5), Inches(5.8), Inches(3.2))
    _fill(u, USER_BG)
    u.line.color.rgb = USER_ACCENT
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(1.75), Inches(1.3), Inches(0.4))
    _fill(badge, USER_ACCENT)
    _add_textbox(slide, Inches(0.95), Inches(1.78), Inches(1.3), Inches(0.35), "你", size=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(0.95), Inches(2.4), Inches(5.2), Inches(1.8),
                 "左侧 / 青色标签\n\n你的提问、纠正、选择方向。\n比如：要筛选 skill、要方案 4、要 MVP 结构、纠正“不是让我自己开发官方 skill”。",
                 size=14, color=INK)
    # ai card
    a = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.8), Inches(3.2))
    _fill(a, AI_BG)
    a.line.color.rgb = AI_ACCENT
    badge2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.05), Inches(1.75), Inches(1.6), Inches(0.4))
    _fill(badge2, AI_ACCENT)
    _add_textbox(slide, Inches(7.05), Inches(1.78), Inches(1.6), Inches(0.35), "Auto", size=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(7.05), Inches(2.4), Inches(5.2), Inches(1.8),
                 "右侧 / 赭石标签\n\n我的回答与澄清。\n尽量保留大白话，不把排行榜和产品混为一谈。",
                 size=14, color=INK)
    _add_textbox(slide, Inches(0.7), Inches(5.1), Inches(11.5), Inches(1.2),
                 "后面按时间顺序回放。长回答会拆成多页，避免一页挤成墙。", size=15, color=MUTED)
    add_footer(slide, prs, "04")


def dialogue_slide(prs, turn: str, speaker: str, title: str, body: str, note: str = ""):
    slide = blank_slide(prs)
    is_user = speaker == "你"
    accent = USER_ACCENT if is_user else AI_ACCENT
    panel_bg = USER_BG if is_user else AI_BG

    _add_textbox(slide, Inches(0.7), Inches(0.35), Inches(8), Inches(0.3), f"对话回放  ·  {turn}", size=12, bold=True, color=MUTED)

    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(0.75), Inches(1.35), Inches(0.38))
    _fill(badge, accent)
    _add_textbox(slide, Inches(0.7), Inches(0.78), Inches(1.35), Inches(0.32), speaker, size=14, bold=True,
                 color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)

    _add_textbox(slide, Inches(2.2), Inches(0.75), Inches(9.5), Inches(0.4), title, size=18, bold=True, color=INK)

    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.35), Inches(11.3), Inches(5.2))
    _fill(panel, panel_bg)
    panel.line.color.rgb = accent

    _add_textbox(slide, Inches(1.0), Inches(1.6), Inches(10.7), Inches(4.5), body, size=15, color=INK)
    if note:
        _add_textbox(slide, Inches(0.7), Inches(6.7), Inches(11.3), Inches(0.3), note, size=11, color=MUTED)
    add_footer(slide, prs, turn)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    cover(prs)
    section(prs, "PART A", "先说清楚方向", "后面那些 AI 对比，不是闲聊，是被需求逼出来的。")
    summary_why(prs)
    summary_plain(prs)
    summary_verdict(prs)
    legend(prs)

    section(prs, "PART B", "完整对话回放", "按发生顺序。你 = 青色；Auto = 赭石。")

    # Dialogue turns condensed but faithful
    dialogues = [
        ("T01", "你", "怎么搜索 skill 并添加？我告诉你，你能筛出来吗？",
         "直接问：能不能按目标帮我筛选现有 skill，并说明该先用哪个。"),
        ("T01", "Auto", "可以。你把场景或关键词给我，我来筛。",
         "当时环境里实际可用的只有：\n"
         "• env-setup（环境）\n"
         "• migrate-to-builds（环境构建迁移）\n"
         "• subscribe（订阅 CI/事件）\n"
         "• walkthrough-artifacts（截图录屏证明）\n\n"
         "并建议你用「需求列表 → 推荐 skill」的格式继续。"),
        ("T02", "你", "想搜：画 CAD、自动处理文档、邮件、自动回微信的 skill",
         "目标很明确：办公 + 工程自动化能力包。"),
        ("T02", "Auto", "按当前可用 skills：没有直接匹配。",
         "结论：CAD / 文档 / 邮件 / 微信 这四类，现成 skill 都未命中。\n"
         "但可以帮你把需求整理成 skill 候选清单，并按落地难度排序：\n"
         "文档 → 邮件 → 微信 → CAD。"),
        ("T03", "你", "好",
         "同意继续整理候选清单。"),
        ("T03", "Auto", "给出四类 skill 候选包",
         "document-automation / email-automation /\n"
         "wechat-auto-reply / cad-drawing-automation\n\n"
         "并说明各自输入输出、依赖、优先级。\n"
         "下一步提供 5 个展开选项。"),
        ("T04", "你", "选 4",
         "要「文档 + 邮件 + 微信」自动化方案，而不是先做 CAD。"),
        ("T04", "Auto", "给出文档+邮件+微信整体方案",
         "五模块：邮件接入、文档处理、规则引擎、消息通知、存储。\n"
         "并拆成三个 skill：文档 / 邮件 / 微信。\n"
         "建议 MVP：读邮件 → 提附件 → 抽字段 → 存表 → 企业微信通知。"),
        ("T05", "你", "再选 3",
         "要 MVP 项目结构设计。"),
        ("T05", "Auto", "给出 Python MVP 目录与模块职责",
         "email / documents / workflows / notifications / storage…\n"
         "强调先打通主链路，别一上来做复杂微信多轮对话。"),
        ("T06", "你", "你是让我自己开发 skill？",
         "关键澄清点：是不是要自己造官方 skill。"),
        ("T06", "Auto", "不是让你开发 Cursor 官方 skill。",
         "两层意思要分开：\n"
         "1）当前环境里的内置 skill：没有你要的那些，搜不到就装不上。\n"
         "2）前面列的 document/email/wechat/cad：是业务能力模块设计，\n"
         "   本质是你自己的自动化系统 / 工作流，不是插件市场一键安装。"),
        ("T07", "你", "Codex 里的 skill 是不是更多？靠那个平台能不能实现？",
         "想确认换平台是否就能直接拥有能力。"),
        ("T07", "Auto", "Codex 生态通常更完整，但不等于自动拥有业务能力。",
         "skill 多是说明书；真要邮件/微信/CAD，还得 MCP/API/脚本。\n"
         "一句话：Codex 更像装配平台，不是现成办公机器人商店。"),
        ("T08", "你", "那就是说装配能力不比 Codex 差，只是说明书少，要我一点点打磨？",
         "开始形成自己的理解模型。"),
        ("T08", "Auto", "基本正确，再补一句：软件本身也重要。",
         "执行和装配能力接近；Codex 往往现成 skill/插件更多。\n"
         "在 Cursor 这边，很多业务能力要自己沉淀成 skill、规则、脚本和工作流。\n"
         "说明书是加速器，不是能力本身。"),
        ("T09", "你", "API 是什么？Auto 和 Codex（ChatGPT）差距有么？",
         "要通俗解释。"),
        ("T09", "Auto", "API = 让两个软件互相说话的接口。",
         "餐厅传菜窗口比喻：你按规则点，厨房按规则上菜。\n"
         "邮箱 API、企业微信 API、CAD API 都是这个意思。\n\n"
         "Auto vs Codex：不完全是智商差距，更像两套工作台。\n"
         "Auto 偏多模型调度；Codex 偏 GPT/OpenAI 生态。"),
        ("T10", "你", "Auto 里有 Kimi、Grok；Codex 是 GPT。是不是像 IBM 集成 vs 苹果自家？好坏只看模型？",
         "半对半错的类比。"),
        ("T10", "Auto", "模型像发动机，平台像整车。",
         "对的部分：Auto 多模型，Codex 更偏 GPT 体系。\n"
         "要修正：好坏不只取决于模型。\n"
         "软件还决定：看得见多少项目、能不能跑命令、有没有连接器、skill 怎么加载。"),
        ("T11", "你", "那为什么 AI 能力排名你比较靠后，写代码最强是 Claude？",
         "把排行榜和工作模式混在一起了。"),
        ("T11", "Auto", "排行榜多测单模型；我是 Auto 调度模式。",
         "Claude 在很多写代码评测里确实常靠前。\n"
         "但 Auto 不是“一个固定大脑”，还带着改仓库、跑命令、接工具的工作方式。\n"
         "排行榜偏裸分，不等于整车体验总分。"),
        ("T12", "你", "好（继续：何时追最强模型）",
         "接受“按任务选模型”的建议。"),
        ("T12", "Auto", "日常推进用 Auto；又难又容易写错时再上最强写代码模型。",
         "对文档/邮件/微信这类自动化：先把流程和接口做出来，比换模型更重要。"),
        ("T13", "你", "阀门定位器程序需要用到 Claude 吗？",
         "落到自己的真实项目。"),
        ("T13", "Auto", "现阶段不必专门上 Claude。",
         "界面、选型、手册公式落地、演示：Auto 通常够。\n"
         "控制回路不稳、固件时序、公式连环错、大重构卡住：再换更强模型。"),
        ("T14", "你", "你这里有 Opus 吗？",
         "确认有没有更强写代码模型可选。"),
        ("T14", "Auto", "软件里有 Opus，但这场对话用的不是它。",
         "现在是 Auto。需要时你要自己切到 Claude Opus。"),
        ("T15", "你", "可以接着说我刚才拦着你、没让你说的问题了",
         "回到：模型 / 平台 / skill / API / 代码各自负责什么。"),
        ("T15", "Auto", "五者分工",
         "模型=脑子；平台=身体工具箱；skill=说明书；\n"
         "API=对外接口；代码=把流程真正跑起来。\n"
         "只有 skill 没有 API 和代码，很多事做不成。"),
        ("T16", "你", "先打断……同样的 skill、同样的模型，放到不同平台，执行同样命令，结果会不一样吗？",
         "精确追问。"),
        ("T16", "Auto", "不一定，但经常会不一样。",
         "纯问答可能接近；一旦依赖项目上下文和工具，平台差异就会放大。\n"
         "同样的脑子，配了不同身体、工具箱、视野和权限。"),
        ("T17", "你", "FY301 动画时，装 Blender 一步步喂你，是不是 API？一步步沟通是不是在提取 skill？",
         "用真实经历验证概念。"),
        ("T17", "Auto", "不完全是 API；喂养非常接近在提取 skill。",
         "装 Blender = 补工具本体；\n"
         "用脚本驱动 Blender = 更像 API；\n"
         "你一步步纠正镜头和讲解顺序 = 沉淀 SOP / skill。\n"
         "效果变好，是因为工具 + 说明书 + 执行路径齐了。"),
        ("T18", "你", "孔板算完要出 CAD 加工图，做不到实际效果；看到 Codex 有人能实现，我们怎么达成？",
         "工程落地痛点。"),
        ("T18", "Auto", "缺的是参数桥 + 模板 + 参考图锁风格 + 稳定出图链路。",
         "不是换到 Codex 就自动会。\n"
         "计算对了，不等于车间认的生产图。\n"
         "正确目标：算完套标准孔板模板，生成可打开 DXF，再 CAD 微调。"),
        ("T19", "你", "我给了国外文丘里图纸模型，也没做出那种效果",
         "补充证据。"),
        ("T19", "Auto", "有参考 ≠ 能复刻。",
         "参考图要拆成可编辑母版、驱动参数表、自动桥、增量规则。\n"
         "文丘里比孔板更复杂，更容易只做出示意。"),
        ("T20", "你", "纠正：文丘里图是要孔板做成那种效果；Codex+skill 是操作 CAD 自己画；问题是风格不像 CAD 生成的，不是尺寸不准。",
         "关键纠正，方向被摆正。"),
        ("T20", "Auto", "承认路线选错了。",
         "你要的是：操作已安装 CAD，从 0 出生产级观感。\n"
         "当时走的是 ezdxf 脚本生成 DXF，容易“一眼机器味”。\n"
         "观感问题先于尺寸问题。正确路径是 CAD 驱动 / 母版 / 操作 CAD 的 skill。"),
    ]

    for turn, speaker, title, body in dialogues:
        dialogue_slide(prs, turn, speaker, title, body)

    # closing
    section(prs, "PART C", "收束", "一句话带走。")
    slide = blank_slide(prs)
    _add_textbox(slide, Inches(0.7), Inches(0.5), Inches(11), Inches(0.4), "带走这句话就够了", size=14, bold=True, color=ACCENT)
    _add_textbox(slide, Inches(0.7), Inches(1.3), Inches(11.5), Inches(2.2),
                 "你要的不是“搜到一个神奇 skill”，\n而是把说明书、接口、代码和平台配齐。",
                 size=26, bold=True, color=INK)
    _add_textbox(slide, Inches(0.7), Inches(3.8), Inches(11.5), Inches(2.0),
                 "后面聊 Auto / Codex / Claude / Opus，\n"
                 "都是为了回答同一件事：\n"
                 "怎么把文档、邮件、微信、CAD、阀门定位器这些活，真正干成。",
                 size=16, color=MUTED)
    add_footer(slide, prs, "尾页")

    ROOT.mkdir(parents=True, exist_ok=True)
    ART.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    prs.save(ART / OUT.name)
    print(f"OK: {OUT}")
    print(f"OK: {ART / OUT.name}")


if __name__ == "__main__":
    build()
